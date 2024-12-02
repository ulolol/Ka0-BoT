#!/usr/bin/env python3

"""
    Ka0$BoT
    Simple WebUI Chatbot with Web Search and Docs Ingestion

    Author: Vidish
"""
import os
import asyncio
import json
import gradio as gr
import requests
from bs4 import BeautifulSoup
import toml
import time
import aiohttp
from typing import List, Dict

class OpenAIWebWrapper:
    def __init__(self, model_name: str = "gpt-4o-mini"):
        self.model_name = model_name
        self.api_key = self.load_api_key()

    def load_api_key(self) -> str:
        """Load the OpenAI API key from a TOML file."""
        try:
            config = toml.load("SearchShellGPT.toml")
            return config['openai']['api_key']
        except Exception as e:
            raise RuntimeError("Failed to load API key from config.toml") from e

    def search_web(self, query: str, num_results: int = 3) -> List[Dict]:
        """Search the web using Google and return results."""
        try:
            #Add more robust import to handle potential package issues
            try:
                from googlesearch import search
            except ImportError:
                print("googlesearch-python package is not installed. Please install it using:")
                print("pip install googlesearch-python")
                return []

            results = []
            try:
                #Add timeout to get around this google blocking issue
                for url in search(query, num_results=num_results, sleep_interval=2):
                    print(f"Checking URL: {url}")
                    results.append({
                        'href': url,
                        'title': self._get_page_title(url),
                        'body': ''
                    })
            except Exception as search_error:
                print(f"Error during web search: {search_error}")
                print("Possible reasons:")
                print("1. No internet connection")
                print("2. Google blocking the search request")
                print("3. googlesearch-python package issues")
                return []

            if not results:
                print(f"No results found for query: {query}")
            
            return results

        except Exception as e:
            print(f"Unexpected error in web search: {e}")
            return []

    def _get_page_title(self, url: str) -> str:
        """Extract the title from a webpage."""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(url, headers=headers, timeout=5)
            response.raise_for_status()

            soup = BeautifulSoup(response.text, 'html.parser')
            title = soup.title.string if soup.title else url
            return title.strip()[:100]  #Limit title length
        except requests.exceptions.RequestException as e:
            print(f"Error fetching title for {url}: {e}")
            return url
        except Exception as e:
            print(f"Unexpected error getting title for {url}: {e}")
            return url
            
    def extract_content(self, url: str) -> str:
        """Extract main content from a webpage."""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'DNT': '1',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1'
            }
            
            response = requests.get(url, headers=headers, timeout=(5, 10))
            
            #Check for specific HTTP errors and handle them gracefully
            if response.status_code == 403:
                print(f"Access forbidden for URL: {url}")
                return ""
            elif response.status_code == 404:
                print(f"URL not found: {url}")
                return ""
            
            response.raise_for_status()

            soup = BeautifulSoup(response.text, 'html.parser')

            for element in soup(['script', 'style', 'nav', 'header', 'footer', 'iframe', 'noscript']):
                element.decompose()

            main_content = soup.find('main') or soup.find('article') or soup.find('div', {'class': ['content', 'main', 'article']})
            if main_content:
                text = main_content.get_text(separator='\n', strip=True)
            else:
                text = soup.get_text(separator='\n', strip=True)

            lines = [line.strip() for line in text.split('\n') if line.strip()]
            cleaned_text = '\n'.join(lines)

            return cleaned_text[:2000]
        
        except requests.exceptions.Timeout:
            print(f"Timeout occurred while extracting content from {url}")
            return ""
        
        except requests.exceptions.ConnectionError:
            print(f"Connection error occurred while extracting content from {url}")
            return ""
        
        except requests.exceptions.HTTPError as e:
            print(f"HTTP error occurred while extracting content from {url}: {e}")
            return ""
        
        except Exception as e:
            print(f"Unexpected error extracting content from {url}: {e}")
            return ""

    def generate_context(self, query: str, num_results: int = 3) -> str:
        """Generate context from web search results."""
        results = self.search_web(query, num_results)
        context = []

        for result in results:
            url = result['href']
            title = result['title']

            try:
                content = self.extract_content(url)

                if content:
                    context_entry = (
                        f"Source: {title}\n"
                        f"URL: {url}\n\n"
                        f"Content:\n{content}\n"
                        f"{'='*50}\n"
                    )
                    context.append(context_entry)
                else:
                    print(f"No content extracted from {url}")
                    context_entry = (
                        f"Source: {title}\n"
                        f"URL: {url}\n"
                        f"No content could be extracted.\n"
                        f"{'='*50}\n"
                    )
                    context.append(context_entry)

            except Exception as e:
                print(f"Error extracting content from {url}: {e}")
                context_entry = (
                    f"Source: {title}\n"
                    f"URL: {url}\n"
                    f"Error extracting content: {e}\n"
                    f"{'='*50}\n"
                )
                context.append(context_entry)

        return "\n".join(context)

    async def query_openai_async(self, prompt: str, context: str, chat_history: list = []) -> any:
        """Query OpenAI with the given prompt and context using aiohttp."""
        try:
            messages = [{"role": "system", "content": "You are a helpful assistant."}]

            #Add chat history to messages
            messages.append({"role": "system", "content": f"Previous chat history for Reference: \n {chat_history}"})
            #messages.extend(chat_history)
            if context:
                messages.append({"role": "system", "content": f"Web Search Context: \n{context}"})

            #Add current prompt and context to messages
            full_prompt = (
                f"Provide a comprehensive answer based on the context above."
                f"Converse with the User naturally." 
                f"Always try to incorporate emojis into your responses naturally and wherever possible."
                f"Always try to keep your responses interesting by using emojis."               
                f"If the user asks a search query, then please provide a comprehensive answer based on the context below, \
                    if available, otherwise use your own knowledge. "
                f"Please ensure the answer is detailed with points wherever necessary. "
                f"Please ensure that the answer is properly formatted for reading. "
                f"If the context doesn't contain relevant information, please state that clearly."
                f"Always format your responses beautifully"
                f"Question: {prompt}\n\n")

            messages.append({"role": "user", "content": full_prompt})

            async with aiohttp.ClientSession() as session:
                async with session.post(
                    'https://api.openai.com/v1/chat/completions',
                    headers={
                        'Authorization': f'Bearer {self.api_key}',
                        'Content-Type': 'application/json'
                    },
                    json={
                        "model": self.model_name,
                        "messages": messages,
                        "max_tokens": 2048,
                        "stream": True
                    }
                ) as response:
                    response_content = ""
                    async for line in response.content:
                        try:
                            line_content = line.decode('utf-8').strip()
                            if not line_content or line_content == "[DONE]":
                                continue
                            
                            if line_content.startswith("data: "):
                                line_content = line_content[len("data: "):]
                            
                            #Skip lines that don't look like JSON
                            if not line_content.startswith('{'):
                                continue
                            
                            json_data = json.loads(line_content)
                            
                            #Check if the response contains a delta with content
                            if 'choices' in json_data and json_data['choices']:
                                delta_content = json_data['choices'][0].get('delta', {}).get('content', '')
                                if delta_content:
                                    response_content += delta_content
                                    yield delta_content
                        except json.JSONDecodeError:
                            #Silently ignore malformed JSON lines
                            continue
                        except Exception as parse_error:
                            #Log or handle any parsing errors
                            print(f"Error parsing response line: {parse_error}")
                            continue

        except Exception as e:
            #More specific error handling
            print(f"Error in query_openai_async: {e}")
            yield f"Error querying OpenAI: {e}"

    async def generate_context_async(self, query: str, num_results: int = 3) -> str:
        """Asynchronously generate context from web search results."""
        results = self.search_web(query, num_results)
        context = []

        async def _async_extract_content(self, url: str) -> str:
            """Asynchronous wrapper for extract_content method"""
            import asyncio
            loop = asyncio.get_event_loop()
            return await loop.run_in_executor(None, self.extract_content, url)

        async def safe_extract_content(result):
            url = result['href']
            title = result['title']

            try:
                #Set a timeout for content extraction
                content = await asyncio.wait_for(
                    _async_extract_content(self, url), 
                    timeout=10.0  #10 seconds timeout
                )

                if content:
                    return (
                        f"Source: {title}\n"
                        f"URL: {url}\n\n"
                        f"Content:\n{content}\n"
                        f"{'='*50}\n"
                    )
                else:
                    return (
                        f"Source: {title}\n"
                        f"URL: {url}\n"
                        f"No content could be extracted.\n"
                        f"{'='*50}\n"
                    )
            except asyncio.TimeoutError:
                print(f"Timeout extracting content from {url}")
                return (
                    f"Source: {title}\n"
                    f"URL: {url}\n"
                    f"Content extraction timed out.\n"
                    f"{'='*50}\n"
                )
            except Exception as e:
                print(f"Error extracting content from {url}: {e}")
                return (
                    f"Source: {title}\n"
                    f"URL: {url}\n"
                    f"Error extracting content: {e}\n"
                    f"{'='*50}\n"
                )

        #Use asyncio to process results concurrently
        try:
            context_results = await asyncio.gather(
                *[safe_extract_content(result) for result in results],
                return_exceptions=True
            )
            
            #Filter out any exceptions and join results
            context = [
                result for result in context_results 
                if not isinstance(result, Exception)
            ]
            
            return "\n".join(context)
        
        except Exception as e:
            print(f"Unexpected error in context generation: {e}")
            return f"Error generating context: {e}"

class GradioAIChatbot:
    def __init__(self):
        self.wrapper = OpenAIWebWrapper(model_name="gpt-4o-mini")
        self.chat_history = []
        self.num_results = 5

    async def chat_response(self, message, history, web_search, file_upload):
        """Generate AI response with optional web search and file context"""
        try:
            #Convert history to OpenAI message format if needed
            formatted_history = [
                {"role": "user" if i % 2 == 0 else "assistant", "content": msg}
                for i, msg in enumerate(sum(history, []))
            ]

            self.chat_history = formatted_history + [{"role": "user", "content": message}]
            
            #Context generation
            context = ""
            
            #Web Search
            if web_search:
                try:
                    context += await self.wrapper.generate_context_async(message, self.num_results)
                except Exception as search_error:
                    print(f"Web search error: {search_error}")
                    context += f"Note: Web search encountered an issue: {search_error}\n"
            
            #File Upload Processing
            if file_upload:
                try:
                    file_content = self.process_file(file_upload.name)
                    context += f"\n\nFile Content:\n{file_content}"
                except Exception as file_error:
                    print(f"File processing error: {file_error}")
                    context += f"Note: File processing encountered an issue: {file_error}\n"

            #Collect streaming response
            response_content = ""
            async_generator = self.wrapper.query_openai_async(message, context, self.chat_history)
            
            try:
                async for chunk in async_generator:
                    response_content += chunk
                    yield response_content
            except Exception as gen_error:
                error_message = f"Error generating response: {gen_error}"
                print(error_message)
                yield error_message

            #Add full response to chat history if successful
            if response_content:
                self.chat_history.append({"role": "assistant", "content": response_content})

        except Exception as overall_error:
            error_message = f"Unexpected error in chat response: {overall_error}"
            print(error_message)
            yield error_message


    def process_file(self, file_path):
        """Process various file types and return content"""
        if not file_path:
            return "No file provided."
        
        #Using the file reading method from the original implementation
        try:
            extension = os.path.splitext(file_path)[1].lower()

            if extension == '.csv':
                import pandas as pd
                df = pd.read_csv(file_path)
                return df.to_string()

            elif extension == '.json':
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return json.dumps(data, indent=2)

            elif extension == '.html':
                from bs4 import BeautifulSoup
                with open(file_path, 'r', encoding='utf-8') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                return soup.get_text()

            elif extension in ('.xlsx', '.xls'):
                import pandas as pd
                df = pd.read_excel(file_path)
                return df.to_string()

            elif extension == '.pdf':
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ''.join(page.extract_text() for page in reader.pages)
                return text

            elif extension == '.docx':
                from docx import Document
                doc = Document(file_path)
                return '\n'.join(para.text for para in doc.paragraphs)

            elif extension == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return '\n'.join(text)

            else:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()

        except Exception as e:
            return f"Error processing file: {e}"


    def launch_gradio(self):
        """Launch the Gradio interface"""
        demo = gr.ChatInterface(
            fn=self.chat_response,
            type='tuples',
            title="🤖 Ka0$BoT",
            description="Chat with an AI assistant that can perform web searches and analyze uploaded files!",
            additional_inputs=[
                gr.Checkbox(label="Enable Web Search"),
                gr.File(label="Upload Context File (optional)")
            ]
        )

        demo.launch()

def main():
    chatbot = GradioAIChatbot()
    chatbot.launch_gradio()

if __name__ == "__main__":
    main()